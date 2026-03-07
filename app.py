from flask import Flask, request, jsonify, send_file
from flask import send_from_directory
import os, json, copy, re, shutil, tempfile, urllib.request
from pptx import Presentation
from lxml import etree
import io, zipfile
from datetime import date, timedelta
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
app = Flask(__name__, static_folder='static')
CONFIG_FILE     = os.path.join(os.path.dirname(__file__), 'config.json')
BIBLE_CACHE_DIR = os.path.join(os.path.dirname(__file__), 'static', 'bible_cache')
BIBLE_META_FILE = os.path.join(os.path.dirname(__file__), 'static', 'bible_meta.json')
os.makedirs(BIBLE_CACHE_DIR, exist_ok=True)

BIBLE_JSON_FILE = os.path.join(BIBLE_CACHE_DIR, "bible_krv.json")
_bible_data_cache = None
_bible_meta_cache = None

# ── Config ────────────────────────────────────

def load_config():
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {'hymn_folder': '', 'template_file': '', 'hymn_slots': []}

def save_config(config):
    with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)

# ── Date helper ───────────────────────────────

def next_sunday_filename():
    today = date.today()
    days_until_sunday = (6 - today.weekday()) % 7
    if days_until_sunday == 0:
        days_until_sunday = 7
    sunday = today + timedelta(days=days_until_sunday)
    return sunday.strftime('%d.%m.%Y') + '.pptx'

# ── Namespace helper ──────────────────────────

_NSMAP = {
    'a':  'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p':  'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r':  'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'pr': 'http://schemas.openxmlformats.org/package/2006/relationships',
    'ct': 'http://schemas.openxmlformats.org/package/2006/content-types',
}

def qn(tag):
    if ':' in tag:
        prefix, local = tag.split(':', 1)
        return '{%s}%s' % (_NSMAP[prefix], local)
    return tag

# ── Bible helpers ─────────────────────────────

BOOK_URL_MAP = {
    'GEN':'gen','EXO':'exo','LEV':'lev','NUM':'num','DEU':'deu','JOS':'jos',
    'JDG':'jdg','RUT':'rut','1SA':'1sa','2SA':'2sa','1KI':'1ki','2KI':'2ki',
    '1CH':'1ch','2CH':'2ch','EZR':'ezr','NEH':'neh','EST':'est','JOB':'job',
    'PSA':'psa','PRO':'pro','ECC':'ecc','SNG':'sng','ISA':'isa','JER':'jer',
    'LAM':'lam','EZK':'ezk','DAN':'dan','HOS':'hos','JOL':'jol','AMO':'amo',
    'OBA':'oba','JON':'jon','MIC':'mic','NAH':'nah','HAB':'hab','ZEP':'zep',
    'HAG':'hag','ZEC':'zec','MAL':'mal','MAT':'mat','MRK':'mrk','LUK':'luk',
    'JHN':'jhn','ACT':'act','ROM':'rom','1CO':'1co','2CO':'2co','GAL':'gal',
    'EPH':'eph','PHP':'php','COL':'col','1TH':'1th','2TH':'2th','1TI':'1ti',
    '2TI':'2ti','TIT':'tit','PHM':'phm','HEB':'heb','JAS':'jas','1PE':'1pe',
    '2PE':'2pe','1JN':'1jn','2JN':'2jn','3JN':'3jn','JUD':'jud','REV':'rev',
}

def load_bible_meta():
    global _bible_meta_cache
    if _bible_meta_cache is None:
        with open(BIBLE_META_FILE, 'r', encoding='utf-8') as f:
            _bible_meta_cache = json.load(f)
    return _bible_meta_cache

def get_bible_chapter(book_id: str, chapter: int) -> dict:
    global _bible_data_cache
    if _bible_data_cache is None:
        if not os.path.exists(BIBLE_JSON_FILE):
            return {"error": f"{BIBLE_JSON_FILE} 파일이 없습니다."}
        with open(BIBLE_JSON_FILE, 'r', encoding='utf-8') as f:
            _bible_data_cache = json.load(f)
    book_data = _bible_data_cache.get(book_id)
    if not book_data:
        return {"error": f"{book_id} 책을 찾을 수 없습니다."}
    chapter_data = book_data.get(str(chapter))
    if not chapter_data:
        return {"error": f"{book_id} {chapter}장을 찾을 수 없습니다."}
    return {int(k): v for k, v in chapter_data.items()}

# ── Routes ────────────────────────────────────

@app.route('/')
def index():
    return send_from_directory('static', 'index.html')

@app.route('/api/config', methods=['GET'])
def get_config():
    return jsonify(load_config())

@app.route('/api/config', methods=['POST'])
def set_config():
    save_config(request.json)
    return jsonify({'ok': True})

@app.route('/api/scan-hymns', methods=['POST'])
def scan_hymns():
    folder = request.json.get('folder', '')
    if not folder or not os.path.exists(folder):
        return jsonify({'error': '폴더를 찾을 수 없어요', 'hymns': []})
    hymns = []
    for f in sorted(os.listdir(folder)):
        if f.lower().endswith('.pptx'):
            m = re.match(r'^(\d+)[\s_\-]*(.*?)\.pptx$', f, re.IGNORECASE)
            if m:
                hymns.append({'number': int(m.group(1)), 'title': m.group(2).strip(), 'filename': f})
            else:
                hymns.append({'number': None, 'title': f.replace('.pptx',''), 'filename': f})
    return jsonify({'hymns': hymns})

@app.route('/api/search-hymn', methods=['POST'])
def search_hymn():
    folder = request.json.get('folder', '')
    number = request.json.get('number')
    if not folder or not os.path.exists(folder):
        return jsonify({'found': False, 'error': '폴더 없음'})
    for f in os.listdir(folder):
        m = re.match(r'^(\d+)[\s_\-]*(.*?)\.pptx$', f, re.IGNORECASE)
        if m and int(m.group(1)) == int(number):
            return jsonify({'found': True, 'filename': f, 'title': m.group(2).strip()})
    return jsonify({'found': False})

@app.route('/api/template-info', methods=['POST'])
def template_info():
    template_file = request.json.get('template_file', '')
    if not template_file or not os.path.exists(template_file):
        return jsonify({'error': '템플릿 파일을 찾을 수 없어요'})
    try:
        prs = Presentation(template_file)
        slides_info = []
        for i, slide in enumerate(prs.slides):
            texts = [s.text_frame.text.strip()[:40]
                     for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            slides_info.append({
                'index': i, 'number': i+1,
                'preview_text': ' / '.join(texts[:2]) if texts else f'슬라이드 {i+1}'
            })
        return jsonify({'slides': slides_info, 'total': len(slides_info)})
    except Exception as e:
        return jsonify({'error': str(e)})

# ── Bible API ─────────────────────────────────

@app.route('/api/bible-meta', methods=['GET'])
def api_bible_meta():
    return jsonify(load_bible_meta())

@app.route('/api/bible', methods=['POST'])
def api_bible():
    data        = request.json or {}
    meta        = load_bible_meta()
    abbr        = data.get('book_abbr', '')
    chapter     = int(data.get('chapter', 1))
    verse_start = int(data.get('verse_start', 1))
    verse_end   = int(data.get('verse_end', verse_start))

    book_id   = meta.get('abbr_to_id', {}).get(abbr)
    book_name = abbr
    if not book_id:
        for b in meta.get('books', []):
            if b['name'] == abbr or b['abbr'] == abbr:
                book_id   = b['id']
                book_name = b['name']
                break
    else:
        for b in meta.get('books', []):
            if b['id'] == book_id:
                book_name = b['name']
                break
    if not book_id:
        return jsonify({'error': f'책을 찾을 수 없어요: {abbr}'})

    chapter_data = get_bible_chapter(book_id, chapter)
    if 'error' in chapter_data:
        return jsonify({'error': f'성경 데이터 오류: {chapter_data["error"]}'})

    verses = []
    for v_num in range(verse_start, verse_end + 1):
        text = chapter_data.get(v_num, '')
        if text:
            verses.append({'num': v_num, 'text': text})
    if not verses:
        return jsonify({'error': f'{verse_start}~{verse_end}절을 찾을 수 없어요'})

    ref_range = f'{verse_start}-{verse_end}' if verse_start != verse_end else str(verse_start)
    return jsonify({
        'verses':      verses,
        'ref':         f'{book_name} {chapter}:{ref_range}',
        'book_name':   book_name,
        'chapter':     chapter,
        'verse_start': verse_start,
        'verse_end':   verse_end,
    })

# ── Generate ──────────────────────────────────

@app.route('/api/generate', methods=['POST'])
def generate():
    work_path = None
    try:
        data          = request.json
        config        = load_config()
        template_file = data.get('template_file') or config.get('template_file', '')
        hymn_folder   = data.get('hymn_folder')   or config.get('hymn_folder', '')
        hymn_slots    = data.get('hymn_slots', [])
        choir         = data.get('choir', None)
        scripture     = data.get('scripture', None)
        extra_verses  = data.get('extra_verses', [])

        if not template_file or not os.path.exists(template_file):
            return jsonify({'error': '템플릿 파일을 찾을 수 없어요'})

        work_fd, work_path = tempfile.mkstemp(suffix='.pptx')
        os.close(work_fd)
        shutil.copy2(template_file, work_path)

        # ── 1. Choir lyrics ───────────────────────────────────────────
        choir_inserted = 0
        if choir and not choir.get('skip'):
            title_idx   = choir.get('title_slide_index')
            song_title  = choir.get('song_title', '').strip()
            lyrics_idx  = choir.get('lyrics_slide_index')
            lyrics_text = choir.get('lyrics', '').strip()
            
            # 제목 슬라이드에 곡명 삽입
            if song_title and title_idx is not None:
                prs = Presentation(work_path)
                set_slide_choir_title(prs.slides[title_idx], song_title)
                prs.save(work_path)

            if lyrics_text and lyrics_idx is not None:
                paragraphs = split_lyrics_into_paragraphs(lyrics_text)
                if paragraphs:
                    prs = Presentation(work_path)
                    set_slide_lyrics(prs.slides[lyrics_idx], paragraphs[0])
                    prs.save(work_path)
                    for i, para in enumerate(paragraphs[1:], 1):
                        insert_after = lyrics_idx + i - 1
                        new_path = duplicate_slide_zip(work_path, lyrics_idx, insert_after)
                        os.unlink(work_path); work_path = new_path
                        prs = Presentation(work_path)
                        set_slide_lyrics(prs.slides[insert_after + 1], para)
                        prs.save(work_path)
                        choir_inserted += 1

        # ── 2. Scripture (시작슬라이드 + 구절슬라이드들) ─────────────────
        if scripture and not scripture.get('skip'):
            sc_title_idx  = scripture.get('title_slide_index')   # 0-based, 시작슬라이드
            sc_idx  = sc_title_idx + 1  # 0-based, 시작슬라이드
            sc_lyrics_idx = scripture.get('lyrics_slide_index')  # 0-based, 구절템플릿
            verses        = scripture.get('verses', [])
            book_name     = scripture.get('book_name', '')
            chapter       = scripture.get('chapter', '')
            verse_start   = scripture.get('verse_start', '')
            verse_end     = scripture.get('verse_end', '')

            # 시작 슬라이드: 책이름 + 장절 교체
            if sc_title_idx is not None and book_name:
                prs = Presentation(work_path)
                set_slide_title_scripture(
                    prs.slides[sc_title_idx], book_name, chapter, verse_start, verse_end
                )
                prs.save(work_path)
            sc_inserted = 0

            # 구절 슬라이드들 (2절씩)
            if verses and sc_idx is not None:
                pairs = []
                for i in range(0, len(verses), 2):
                    chunk = verses[i:i+2]
                    pairs.append(chunk)
        
                prs = Presentation(work_path)
                target_slide = prs.slides[sc_idx]

                set_slide_text_bibel(prs.slides[sc_idx], pairs[0])
                add_chapter_title_text(target_slide, f"{book_name} {chapter}장")

                prs.save(work_path)

                for i, pair in enumerate(pairs[1:], 1):
                    insert_after = sc_idx + i - 1
                    new_path = duplicate_slide_zip(work_path, sc_idx, insert_after)
                    os.unlink(work_path); work_path = new_path
                    prs = Presentation(work_path)
                    set_slide_text_bibel(prs.slides[insert_after + 1], pair)
                    target_slide = prs.slides[insert_after + 1]
                    add_chapter_title_text(target_slide, f"{book_name} {chapter}장")

                    prs.save(work_path)
                    sc_inserted += 1
        else:
            sc_inserted = 0


        # ── 3. 추가 구절 (extra_verses) ───────────────────────────────
        ev_in_order = sorted(
            [ev for ev in extra_verses if ev.get('slide_index') is not None and ev.get('verses')],
            key=lambda x: x['slide_index']
        )

        # 역순 처리: 인덱스 큰 것부터 삽입 → 앞쪽 인덱스에 영향 없음
        # delete_slide는 루프 밖에서 일괄 처리 (루프 안 삭제 시 인덱스 틀어짐)
        ev_indices_to_delete = []

        for ev_orig_order, ev in enumerate(reversed(ev_in_order)):
            raw_idx = ev['slide_index']
            sc_lyrics_idx_val = scripture.get('lyrics_slide_index') if scripture and not scripture.get('skip') else None
            choir_lyrics_idx_val = choir.get('lyrics_slide_index') if choir and not choir.get('skip') else None

            ev_idx = raw_idx
            if choir_lyrics_idx_val is not None and raw_idx > choir_lyrics_idx_val:
                ev_idx += choir_inserted
            if sc_lyrics_idx_val is not None and raw_idx > sc_lyrics_idx_val:
                ev_idx += sc_inserted

            ev_verses = ev['verses']
            ev_book = ev.get('book_name', '')
            ev_ch   = str(ev.get('chapter', ''))

            prs_check = Presentation(work_path)
            total = len(prs_check.slides)
            if ev_idx >= total:
                continue

            pairs = [ev_verses[i:i+2] for i in range(0, len(ev_verses), 2)]

            original_order_idx = (len(ev_in_order) - 1) - ev_orig_order
            if original_order_idx >= 1:
                from pptx.util import Pt
                from pptx.dml.color import RGBColor
                from pptx.util import Emu
                import copy

                # 빈 검은 슬라이드 새로 삽입
                prs = Presentation(work_path)
                blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
                new_slide = prs.slides.add_slide(blank_layout)

                # 배경을 검정으로
                from pptx.oxml.ns import qn
                from lxml import etree
                bg = new_slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 0, 0)

                # 삽입 위치: ev_idx+1 로 이동 (add_slide는 맨 뒤에 추가됨 → XML 재정렬)
                xml_slides = prs.slides._sldIdLst
                last = xml_slides[-1]
                xml_slides.remove(last)
                xml_slides.insert(ev_idx + 1, last)

                # 구분 슬라이드가 ev_idx+1 위치에 들어갔으므로
                # 이후 구절 슬라이드들은 ev_idx+1 뒤에 삽입
                prs.save(work_path)
                base = ev_idx + 1
            else:
                base = ev_idx

            for i, pair in enumerate(pairs):
                ins_after = base + i
                new_path = duplicate_slide_zip(work_path, ev_idx, ins_after)
                os.unlink(work_path); work_path = new_path

                prs = Presentation(work_path)
                target_slide = prs.slides[ins_after + 1]
                set_slide_text_bibel(target_slide, pair)
                add_chapter_title_text(target_slide, f"{ev_book} {ev_ch}장")
                prs.save(work_path)

            ev_indices_to_delete.append(ev_idx)

        # 원본 템플릿 슬라이드들 일괄 삭제 (내림차순으로 삭제해야 인덱스 안 밀림)
        if ev_indices_to_delete:
            prs = Presentation(work_path)
            for idx in sorted(ev_indices_to_delete, reverse=True):
                delete_slide(prs, idx)
            prs.save(work_path)


        # ── 4. Hymn insertion (back-to-front) ─────────────────────────
        active = [s for s in hymn_slots if not s.get('skip') and s.get('hymn_number')]
        active = sorted(active, key=lambda x: x.get('after_slide_index', 0), reverse=True)

        for slot in active:
            hymn_file = find_hymn_file(hymn_folder, slot['hymn_number'])
            if not hymn_file:
                continue
            hymn_prs = Presentation(hymn_file)
            n        = len(hymn_prs.slides)
            after    = slot.get('after_slide_index', 0)
            for i in range(n):
                new_path = copy_slide_from_file_zip(hymn_file, i, work_path, after + i)
                os.unlink(work_path); work_path = new_path

        # ── 5. Return ────────────────────────────────────────────────
        filename = next_sunday_filename()
        with open(work_path, 'rb') as f:
            data_bytes = f.read()
        os.unlink(work_path); work_path = None

        return send_file(
            io.BytesIO(data_bytes),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        import traceback
        if work_path and os.path.exists(work_path):
            try: os.unlink(work_path)
            except: pass
        return jsonify({'error': str(e), 'trace': traceback.format_exc()})


# ── Core ZIP-level slide copy ─────────────────

def _read_xml(zf, path):
    with zf.open(path) as f:
        return etree.parse(f).getroot()

def _xml_bytes(root):
    return etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)

def _slide_paths_ordered(zf):
    prs_xml  = _read_xml(zf, 'ppt/presentation.xml')
    prs_rels = _read_xml(zf, 'ppt/_rels/presentation.xml.rels')
    rid_to_t = {r.get('Id'): r.get('Target','') for r in prs_rels}
    paths = []
    sldIdLst = prs_xml.find(qn('p:sldIdLst'))
    if sldIdLst is None:
        return paths
    for sldId in sldIdLst:
        t = rid_to_t.get(sldId.get(qn('r:id')), '')
        if t:
            paths.append('ppt/' + t.lstrip('./'))
    return paths

def _rels_path(slide_path):
    parts = slide_path.rsplit('/', 1)
    return parts[0] + '/_rels/' + parts[1] + '.rels'

def _max_slide_num(zf):
    nums = [int(m.group(1)) for n in zf.namelist()
            for m in [re.match(r'ppt/slides/slide(\d+)\.xml$', n)] if m]
    return max(nums) if nums else 0

def _max_rid(rels_root):
    ids = [int(m.group(1)) for r in rels_root
           for m in [re.match(r'rId(\d+)', r.get('Id',''))] if m]
    return max(ids) + 1 if ids else 1

def _max_sld_id(prs_xml):
    sldIdLst = prs_xml.find(qn('p:sldIdLst'))
    ids = [int(el.get('id', 255)) for el in (list(sldIdLst) if sldIdLst is not None else [])]
    return max(ids) + 1 if ids else 256

EXT_MIME = {
    'png':'image/png','jpg':'image/jpeg','jpeg':'image/jpeg',
    'gif':'image/gif','bmp':'image/bmp','tiff':'image/tiff',
    'svg':'image/svg+xml','wmf':'image/x-wmf','emf':'image/x-emf',
    'mp4':'video/mp4','mp3':'audio/mpeg','wav':'audio/wav',
}

def copy_slide_from_file_zip(src_path, src_slide_index, dst_path, insert_after):
    out_fd, out_path = tempfile.mkstemp(suffix='.pptx')
    os.close(out_fd)

    with zipfile.ZipFile(src_path, 'r') as src_zf, \
         zipfile.ZipFile(dst_path, 'r') as dst_zf:

        src_slide_paths = _slide_paths_ordered(src_zf)
        if src_slide_index >= len(src_slide_paths):
            raise ValueError(f'Slide index {src_slide_index} out of range')
        src_slide_path = src_slide_paths[src_slide_index]
        src_rels_path  = _rels_path(src_slide_path)

        dst_prs_xml  = _read_xml(dst_zf, 'ppt/presentation.xml')
        dst_prs_rels = _read_xml(dst_zf, 'ppt/_rels/presentation.xml.rels')
        dst_ct_xml   = _read_xml(dst_zf, '[Content_Types].xml')

        new_slide_num  = _max_slide_num(dst_zf) + 1
        new_slide_path = f'ppt/slides/slide{new_slide_num}.xml'
        new_rels_path  = f'ppt/slides/_rels/slide{new_slide_num}.xml.rels'
        new_rid        = f'rId{_max_rid(dst_prs_rels)}'
        new_sld_id     = _max_sld_id(dst_prs_xml)

        src_rels_root = etree.Element('{%s}Relationships' % _NSMAP['pr'])
        if src_rels_path in src_zf.namelist():
            src_rels_root = _read_xml(src_zf, src_rels_path)

        dst_names    = set(dst_zf.namelist())
        extra_files  = {}
        ct_defaults  = {el.get('Extension','').lower() for el in dst_ct_xml.findall(qn('ct:Default'))}
        ct_overrides = {el.get('PartName','')           for el in dst_ct_xml.findall(qn('ct:Override'))}

        new_rels_root = etree.Element('{%s}Relationships' % _NSMAP['pr'])

        for rel in src_rels_root:
            rid      = rel.get('Id','')
            rel_type = rel.get('Type','')
            target   = rel.get('Target','')
            tmode    = rel.get('TargetMode','')

            new_rel = etree.SubElement(new_rels_root, qn('pr:Relationship'))
            new_rel.set('Id',   rid)
            new_rel.set('Type', rel_type)

            if tmode == 'External':
                new_rel.set('Target', target)
                new_rel.set('TargetMode', 'External')
                continue

            if target.startswith('../'):
                src_full = 'ppt/' + target[3:]
            else:
                src_full = 'ppt/slides/' + target

            if 'slideLayout' in target or 'slideMaster' in target:
                new_rel.set('Target', target)
                continue

            if src_full not in src_zf.namelist():
                new_rel.set('Target', target)
                continue

            file_bytes = src_zf.read(src_full)
            ext = os.path.splitext(src_full)[1].lower().lstrip('.')

            dst_target = src_full
            counter = 2
            all_used = dst_names | set(extra_files.keys())
            while dst_target in all_used:
                base, dot_ext = os.path.splitext(src_full)
                dst_target = f'{base}_{counter}{dot_ext}'
                counter += 1

            extra_files[dst_target] = file_bytes

            if ext and ext not in ct_defaults:
                mime = EXT_MIME.get(ext, 'application/octet-stream')
                nd = etree.SubElement(dst_ct_xml, qn('ct:Default'))
                nd.set('Extension', ext)
                nd.set('ContentType', mime)
                ct_defaults.add(ext)

            new_rel.set('Target', '../' + dst_target[4:])

        part_name = '/' + new_slide_path
        if part_name not in ct_overrides:
            ov = etree.SubElement(dst_ct_xml, qn('ct:Override'))
            ov.set('PartName', part_name)
            ov.set('ContentType',
                   'application/vnd.openxmlformats-officedocument.presentationml.slide+xml')

        pr_rel = etree.SubElement(dst_prs_rels, qn('pr:Relationship'))
        pr_rel.set('Id',     new_rid)
        pr_rel.set('Type',   'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
        pr_rel.set('Target', f'slides/slide{new_slide_num}.xml')

        sldIdLst = dst_prs_xml.find(qn('p:sldIdLst'))
        new_sldId_el = etree.Element(qn('p:sldId'))
        new_sldId_el.set('id',       str(new_sld_id))
        new_sldId_el.set(qn('r:id'), new_rid)

        children = list(sldIdLst)
        pos = min(insert_after + 1, len(children))
        for c in children: sldIdLst.remove(c)
        children.insert(pos, new_sldId_el)
        for c in children: sldIdLst.append(c)

        with zipfile.ZipFile(out_path, 'w', zipfile.ZIP_DEFLATED) as out_zf:
            skip = {'ppt/presentation.xml', 'ppt/_rels/presentation.xml.rels', '[Content_Types].xml'}
            for name in dst_zf.namelist():
                if name not in skip:
                    out_zf.writestr(name, dst_zf.read(name))
            out_zf.writestr(new_slide_path, src_zf.read(src_slide_path))
            out_zf.writestr(new_rels_path,  _xml_bytes(new_rels_root))
            for dst_name, fb in extra_files.items():
                if dst_name not in dst_zf.namelist():
                    out_zf.writestr(dst_name, fb)
            out_zf.writestr('ppt/presentation.xml',            _xml_bytes(dst_prs_xml))
            out_zf.writestr('ppt/_rels/presentation.xml.rels', _xml_bytes(dst_prs_rels))
            out_zf.writestr('[Content_Types].xml',             _xml_bytes(dst_ct_xml))

    return out_path


def duplicate_slide_zip(pptx_path, slide_index, insert_after):
    return copy_slide_from_file_zip(pptx_path, slide_index, pptx_path, insert_after)


# ── Lyrics helpers ────────────────────────────
def set_slide_choir_title(slide, song_title: str):
    """
    성가대 제목 슬라이드: 가장 큰 텍스트박스의 첫 번째 run 내용만 교체.
    서식(폰트 크기/색상 등)은 유지.
    """
    best_shape, best_size = None, 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            size = shape.width * shape.height
            if size > best_size:
                best_size  = size
                best_shape = shape
    if best_shape is None:
        return

    tf     = best_shape.text_frame
    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))

    for para in paras:
        runs = para.findall(qn('a:r'))
        if runs:
            t_el = runs[0].find(qn('a:t'))
            if t_el is None:
                t_el = etree.SubElement(runs[0], qn('a:t'))
            t_el.text = song_title
            for extra_r in runs[1:]:
                para.remove(extra_r)
            break  # 첫 번째 단락만 교체

def split_lyrics_into_paragraphs(text):
    return [p.strip() for p in re.split(r'\n\s*\n', text.strip()) if p.strip()]


def set_slide_lyrics(slide, lyrics_text):
    best_shape, best_size = None, 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            size = shape.width * shape.height
            if size > best_size:
                best_size  = size
                best_shape = shape
    if best_shape is None:
        return
    tf     = best_shape.text_frame
    lines  = lyrics_text.split('\n')
    txBody = tf._txBody
    existing_ps = txBody.findall(qn('a:p'))
    template_p  = copy.deepcopy(existing_ps[0]) if existing_ps else None
    for p in existing_ps:
        txBody.remove(p)
    for line in lines:
        if template_p is not None:
            new_p = copy.deepcopy(template_p)
            for r in new_p.findall(qn('a:r')):
                new_p.remove(r)
            runs = template_p.findall(qn('a:r'))
            if runs:
                new_r = copy.deepcopy(runs[0])
                t_el  = new_r.find(qn('a:t'))
                if t_el is None:
                    t_el = etree.SubElement(new_r, qn('a:t'))
                t_el.text = line
                new_p.append(new_r)
            else:
                new_r = etree.SubElement(new_p, qn('a:r'))
                t_el  = etree.SubElement(new_r, qn('a:t'))
                t_el.text = line
        else:
            new_p = etree.Element(qn('a:p'))
            new_r = etree.SubElement(new_p, qn('a:r'))
            t_el  = etree.SubElement(new_r, qn('a:t'))
            t_el.text = line
        txBody.append(new_p)


# ── Bible slide helpers ───────────────────────

def set_slide_title_scripture(slide, book_name: str, chapter, verse_start, verse_end):
    """
    성경구절 시작슬라이드: 기존 paragraph/run 구조와 서식(폰트크기 포함)을 그대로 유지,
    텍스트 내용만 교체.
      paragraph 1 (또는 run 1) → 책이름  (예: 요한복음)
      paragraph 2 (또는 run 2) → 장절    (예: 8장 3-9절)
    paragraph가 1개뿐이면 run을 2개로 나눠서 처리.
    """
    best_shape, best_size = None, 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            size = shape.width * shape.height
            if size > best_size:
                best_size  = size
                best_shape = shape
    if best_shape is None:
        return

    verse_range = (f'{verse_start}-{verse_end}절'
                   if str(verse_start) != str(verse_end) else f'{verse_start}절')
    chapter_verse_text = f'{chapter}장 {verse_range}'

    tf     = best_shape.text_frame
    txBody = tf._txBody
    paras  = txBody.findall(qn('a:p'))

    if len(paras) >= 2:
        # paragraph별 첫 run 내용만 교체, 나머지 run 제거 (서식 유지)
        for para, new_text in zip(paras[:2], [book_name, chapter_verse_text]):
            runs = para.findall(qn('a:r'))
            if runs:
                t_el = runs[0].find(qn('a:t'))
                if t_el is None:
                    t_el = etree.SubElement(runs[0], qn('a:t'))
                t_el.text = new_text
                for extra_r in runs[1:]:
                    para.remove(extra_r)
            else:
                new_r = etree.SubElement(para, qn('a:r'))
                t_el  = etree.SubElement(new_r, qn('a:t'))
                t_el.text = new_text
    else:
        # paragraph 1개: 기존 run 서식 복제 → run 2개로
        para = paras[0] if paras else etree.SubElement(txBody, qn('a:p'))
        runs = para.findall(qn('a:r'))
        tmpl_run = copy.deepcopy(runs[0]) if runs else None
        for r in runs:
            para.remove(r)
        for new_text in [book_name, chapter_verse_text]:
            new_r = copy.deepcopy(tmpl_run) if tmpl_run else etree.Element(qn('a:r'))
            t_el  = new_r.find(qn('a:t'))
            if t_el is None:
                t_el = etree.SubElement(new_r, qn('a:t'))
            t_el.text = new_text
            para.append(new_r)


def set_slide_text_bibel(slide, text_content):

    best_shape, best_size = None, 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            size = shape.width * shape.height
            if size > best_size:
                best_size  = size
                best_shape = shape

    if best_shape is None:
        print("성경구절 생성 오류: 택스트 박스를 찾지 못했습니다.")
        return

    tf = best_shape.text_frame
    txBody = tf._txBody
    exist_p = txBody.findall(qn('a:p'))

    # 첫 paragraph run 서식 가져오기
    tmpl_run = None
    if exist_p:
        first_p = exist_p[0]
        runs = first_p.findall(qn('a:r'))
        if runs:
            tmpl_run = copy.deepcopy(runs[0])

    # 기존 paragraph 삭제
    for p in exist_p:
        txBody.remove(p)
        
    tab_pos = 705000  # 번호 뒤 텍스트 시작 위치
    tab_pos_use = 0
    new_lines = []
    
    for verse in text_content:
        # 글자 자동 줄바꿈용 글자 크기
        font_size_pt = 18  # 기본값
        if tmpl_run is not None and hasattr(tmpl_run, 'rPr') and hasattr(tmpl_run.rPr, 'sz'):
            sz = tmpl_run.rPr.sz
            font_size_pt = sz / 100 if isinstance(sz, int) else sz.pt

        wrapped = estimate_line_breaks(verse["text"], best_shape.width, font_size_pt)

        for i, l in enumerate(wrapped):
            new_lines.append({
                "num": verse["num"] if i == 0 else None,  # 첫 줄만 번호
                "text": l
            })
    
    for line in new_lines:

        new_p = etree.Element(qn('a:p'))
        pPr = etree.SubElement(new_p, qn('a:pPr'))

        # tab 위치 결정
        tab_pos_use = tab_pos

        # tab 설정
        tabLst = etree.SubElement(pPr, qn('a:tabLst'))
        tab = etree.SubElement(tabLst, qn('a:tab'))
        tab.set('pos', str(tab_pos_use))

        # run 생성 (서식 유지)
        if tmpl_run is not None:
            new_r = copy.deepcopy(tmpl_run)
        else:
            new_r = etree.SubElement(new_p, qn('a:r'))

        t_el = new_r.find('.//' + qn('a:t'))
        if t_el is None:
            t_el = etree.SubElement(new_r, qn('a:t'))

        # 번호 줄 / 번호 없는 줄 텍스트
        if line["num"] is not None:
            t_el.text = f"{line['num']}.\t{line['text']}"
        else:
            t_el.text = f"\t{line['text']}"

        new_p.append(new_r)
        txBody.append(new_p)

# 한 줄에 들어갈 문자 수 계산 (대략)
def estimate_line_breaks(text, box_width_emu, font_size_pt):

    EMU_PER_PT = 12700

    # 텍스트박스 width를 pt로 변환
    box_width_pt = box_width_emu / EMU_PER_PT

    # 한국어 평균 글자폭 (폰트의 약 0.9배)
    char_width_pt = font_size_pt * 0.9

    max_chars = int(box_width_pt / char_width_pt)

    if max_chars <= 1:
        max_chars = 10

    words = text.split()
    lines = []
    current = ""

    for w in words:

        test = (current + " " + w).strip()

        if len(test) <= max_chars:
            current = test
        else:
            lines.append(current)
            current = w

    if current:
        lines.append(current)

    return lines

def add_chapter_title_text(slide, title_text: str):
    """
    슬라이드 안의 가장 작은 텍스트박스를 찾아 내용을 title_text로 교체.
    글자 서식(rPr)은 기존 첫 번째 run에서 deepcopy하여 재사용.
    텍스트박스가 없으면 새로 추가.
    """
    import copy
    from lxml import etree
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    # 가장 작은 텍스트박스 탐색 (초기값 inf → 작은 것 선택)
    small_shape, small_size = None, float('inf')
    for shape in slide.shapes:
        if shape.has_text_frame:
            size = shape.width * shape.height
            if size < small_size:
                small_size  = size
                small_shape = shape

    if small_shape is not None:
        tf = small_shape.text_frame

        # 기존 첫 번째 run의 rPr(글자 서식) 복제
        source_rPr = None
        for para in tf.paragraphs:
            for run in para.runs:
                rPr_el = run._r.find(f'{{{ns}}}rPr')
                if rPr_el is not None:
                    source_rPr = copy.deepcopy(rPr_el)
                    break
            if source_rPr is not None:
                break

        # 모든 단락의 런 초기화 (텍스트 지우기)
        for para in tf.paragraphs:
            for r in list(para._p.findall(f'{{{ns}}}r')):
                para._p.remove(r)

#         # python-pptx 고수준 API로 런 추가 → XML 구조 보장
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = title_text

        # source_rPr 서식 복제 적용
        if source_rPr is not None:
            run._r.insert(0, source_rPr)
            
def delete_slide(prs, slide_index):
    xml_slides = prs.slides._sldIdLst
    slide_el = xml_slides[slide_index]
    xml_slides.remove(slide_el)

# ── File helper ────────────────────────────────

def find_hymn_file(folder, number):
    if not folder or not os.path.exists(folder):
        return None
    for f in os.listdir(folder):
        m = re.match(r'^(\d+)[\s_\-]*(.*?)\.pptx$', f, re.IGNORECASE)
        if m and int(m.group(1)) == int(number):
            return os.path.join(folder, f)
    return None


if __name__ == '__main__':
    os.makedirs('static', exist_ok=True)
    os.makedirs(BIBLE_CACHE_DIR, exist_ok=True)
    print('✝  교회 PPT 자동화 도구 시작!')
    print('브라우저에서 http://localhost:5000 을 열어주세요')
    app.run(debug=True, port=5000)